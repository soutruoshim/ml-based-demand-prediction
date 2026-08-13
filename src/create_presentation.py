#!/usr/bin/env python3
"""Create the final-defense presentation from verified project outputs."""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs"
DEST = ROOT / "presentation" / "ML_Demand_Prediction_Final_Defense.pptx"
NAVY, TEAL, WHITE, GRAY = RGBColor(18, 52, 77), RGBColor(23, 107, 135), RGBColor(255,255,255), RGBColor(70,78,86)

def style(slide, title):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    shape = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(.72))
    shape.fill.solid(); shape.fill.fore_color.rgb = NAVY; shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(.55), Inches(.12), Inches(12.2), Inches(.45))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(25); p.font.bold=True; p.font.color.rgb=WHITE

def bullets(prs, title, items, note=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); style(s,title)
    box=s.shapes.add_textbox(Inches(.8), Inches(1.15), Inches(11.8), Inches(5.4)); tf=box.text_frame
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.size=Pt(23); p.font.color.rgb=GRAY; p.space_after=Pt(18); p.level=0
    if note:
        n=s.shapes.add_textbox(Inches(.8), Inches(6.65), Inches(11.8), Inches(.4)); p=n.text_frame.paragraphs[0]; p.text=note; p.font.size=Pt(11); p.font.italic=True; p.font.color.rgb=TEAL

def image_slide(prs,title,path,caption):
    s=prs.slides.add_slide(prs.slide_layouts[6]); style(s,title)
    s.shapes.add_picture(str(path), Inches(1.45), Inches(1.02), width=Inches(10.4), height=Inches(5.75))
    b=s.shapes.add_textbox(Inches(.8), Inches(6.85), Inches(11.8), Inches(.35)); p=b.text_frame.paragraphs[0]; p.text=caption; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(12); p.font.color.rgb=GRAY

def main():
    eval_df=pd.read_csv(OUT/"model_evaluation_results.csv"); plan=pd.read_csv(OUT/"jan_2026_production_recommendation.csv")
    best=eval_df.iloc[0]
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
    tb=s.shapes.add_textbox(Inches(.8), Inches(1.55), Inches(11.7), Inches(2.2)); tf=tb.text_frame
    p=tf.paragraphs[0]; p.text="Machine Learning-Based Demand Prediction"; p.font.size=Pt(36); p.font.bold=True; p.font.color.rgb=WHITE; p.alignment=PP_ALIGN.CENTER
    p=tf.add_paragraph(); p.text="Factory Production Planning for Shampoo, Body Wash, and Laundry Detergent"; p.font.size=Pt(24); p.font.color.rgb=RGBColor(100,204,197); p.alignment=PP_ALIGN.CENTER
    p=tf.add_paragraph(); p.text="Group 16 · SOUT RAHIM · Norton University · August 2026"; p.font.size=Pt(18); p.font.color.rgb=WHITE; p.space_before=Pt(32); p.alignment=PP_ALIGN.CENTER
    bullets(prs,"Problem and Goal",["Overproduction raises inventory and storage costs; underproduction risks stockouts and missed sales.","Goal: predict next-month demand and translate the forecast into a clear production recommendation.","Scope: Shampoo, Body Wash, and Laundry Detergent."],"Supervised machine-learning regression project")
    bullets(prs,"Data and Academic Integrity",["177 monthly product records: February 2021–December 2025.","Features: product, month/season, previous sales, stock, price, promotion, customer orders, and previous production.","The dataset is synthetic and is used only to demonstrate the workflow; it is not evidence of actual Medtherm demand."],"Replace the CSV with authorized factory records before operational use.")
    bullets(prs,"Methodology",["Sort observations by month and reserve the latest 20% as a holdout test period (141 train / 36 test).","Fit imputation, one-hot encoding, and scaling only on training data to prevent leakage.","Compare Linear Regression, Decision Tree, and Random Forest on identical data.","Select the lowest-RMSE model; verify with MAE and R²."])
    image_slide(prs,"Model Comparison",OUT/"figures/model_rmse.png",f"Selected model: {best['Model']} · MAE {best['MAE']:.1f} · RMSE {best['RMSE']:.1f} · R² {best['R2']:.3f}")
    image_slide(prs,"Prediction Quality",OUT/"figures/actual_vs_predicted.png","Held-out observations near the diagonal indicate closer predictions.")
    image_slide(prs,"Demand Trends",OUT/"figures/demand_trends.png","Actual and predicted demand across the final 12 held-out months.")
    image_slide(prs,"What Drives the Forecast?",OUT/"figures/feature_importance.png","Permutation importance measures the increase in error when each feature is shuffled.")
    bullets(prs,"Production Planning Rule",["Recommended production = max(0, predicted demand + 10% safety stock − current stock).",*(f"{r.product_type}: predict {r.predicted_demand:,.0f}; recommend {r.recommended_production:,.0f} units" for r in plan.itertuples())],"Illustrative January 2026 scenario; the safety factor is a demonstration policy, not an optimized factory rule.")
    bullets(prs,"Practical Value",["Managers receive a repeatable forecast instead of relying only on experience or simple averages.","Production teams can prepare materials and schedules earlier.","Sales and inventory teams can review demand, stock, and orders in one planning output.","Human approval remains essential because capacity, lead time, and batch constraints are not modeled."])
    bullets(prs,"Limitations and Future Work",["Synthetic, small dataset with only three product categories.","External drivers such as holidays, channels, competitor activity, and weather are absent.","Next: use authorized factory data, time-series cross-validation, tuned models, prediction intervals, and capacity-aware optimization."])
    bullets(prs,"Conclusion",[f"A complete, reproducible regression workflow was implemented and evaluated.",f"{best['Model']} achieved the best synthetic holdout RMSE ({best['RMSE']:.1f} units).","The project demonstrates how forecasts can support production decisions while clearly separating demonstration evidence from real factory claims.","Questions?"],"All code, data, outputs, figures, and the paper are included in the submission folder.")
    DEST.parent.mkdir(exist_ok=True); prs.save(DEST); print(DEST)

if __name__ == "__main__": main()
